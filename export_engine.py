"""
export_engine.py
----------------
Generates PDF and PowerPoint exports of financial analysis.
Both formats are professional, dark-themed, and FP&A-ready.

PDF:  reportlab — structured CFO brief with KPIs, health score, insights
PPTX: python-pptx — 6-slide deck with charts and key metrics
"""

import io
import datetime
import pandas as pd
import numpy as np


# ─── Colour palette (matches app) ────────────────────────────────────────────
C_BG       = (0,   0,   0)        # #000000
C_CARD     = (28,  28,  30)       # #1C1C1E
C_BORDER   = (44,  44,  46)       # #2C2C2E
C_WHITE    = (255, 255, 255)
C_MUTED    = (142, 142, 147)      # #8E8E93
C_BLUE     = (10,  132, 255)      # #0A84FF
C_GREEN    = (52,  199, 89)       # #34C759
C_RED      = (255, 59,  48)       # #FF3B30
C_ORANGE   = (255, 159, 10)       # #FF9F0A
C_PURPLE   = (191, 90,  242)      # #BF5AF2


def _rgb(c):
    """Convert tuple to 0-1 float triple for reportlab."""
    return tuple(v/255 for v in c)


# ══════════════════════════════════════════════════════════════════════════════
# PDF EXPORT
# ══════════════════════════════════════════════════════════════════════════════

def generate_pdf(
    company_name: str,
    ticker: str,
    info: dict,
    kpis: dict,
    health_score: int,
    health_label: str,
    health_breakdown: dict,
    insights: dict,
    cfo_brief: str,
    news_items: list[dict],
) -> bytes:
    """Generate a professional PDF financial brief. Returns bytes."""
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import mm
        from reportlab.lib import colors
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Table,
            TableStyle, HRFlowable, PageBreak,
        )
        from reportlab.lib.styles import ParagraphStyle
        from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT

        buf = io.BytesIO()
        W, H = A4

        # Custom styles
        def style(name, **kw):
            defaults = dict(
                fontName="Helvetica", fontSize=10,
                textColor=colors.Color(*_rgb(C_WHITE)),
                backColor=colors.Color(*_rgb(C_BG)),
                leading=14, spaceAfter=4,
            )
            defaults.update(kw)
            return ParagraphStyle(name, **defaults)

        s_title   = style("title",   fontSize=22, fontName="Helvetica-Bold", spaceAfter=4)
        s_sub     = style("sub",     fontSize=11, textColor=colors.Color(*_rgb(C_MUTED)))
        s_h1      = style("h1",      fontSize=14, fontName="Helvetica-Bold",
                          textColor=colors.Color(*_rgb(C_BLUE)), spaceBefore=12, spaceAfter=6)
        s_h2      = style("h2",      fontSize=11, fontName="Helvetica-Bold", spaceBefore=8, spaceAfter=4)
        s_body    = style("body",    fontSize=9,  leading=14,
                          textColor=colors.Color(*_rgb(C_WHITE)))
        s_muted   = style("muted",   fontSize=8,
                          textColor=colors.Color(*_rgb(C_MUTED)))
        s_label   = style("label",   fontSize=8,  fontName="Helvetica-Bold",
                          textColor=colors.Color(*_rgb(C_MUTED)))

        label_color_map = {
            "Excellent": colors.Color(*_rgb(C_GREEN)),
            "Strong":    colors.Color(*_rgb(C_BLUE)),
            "Average":   colors.Color(*_rgb(C_ORANGE)),
            "Weak":      colors.Color(*_rgb(C_RED)),
        }
        score_color = label_color_map.get(health_label, colors.Color(*_rgb(C_MUTED)))

        story = []
        margin = 18 * mm
        usable_w = W - 2 * margin

        # ── Cover header ─────────────────────────────────────────────────────
        today = datetime.date.today().strftime("%B %d, %Y")
        story.append(Paragraph(f"FinIntel AI — Financial Intelligence Report", s_muted))
        story.append(Spacer(1, 4))
        story.append(Paragraph(company_name, s_title))
        story.append(Paragraph(
            f"{ticker} &nbsp;·&nbsp; {info.get('sector','N/A')} &nbsp;·&nbsp; "
            f"{info.get('country','N/A')} &nbsp;·&nbsp; Generated {today}",
            s_sub,
        ))
        story.append(HRFlowable(width="100%", thickness=1,
                                color=colors.Color(*_rgb(C_BORDER)), spaceAfter=12))

        # ── Health score strip ────────────────────────────────────────────────
        score_data = [[
            Paragraph("Financial Health Score", s_label),
            Paragraph("Market Cap", s_label),
            Paragraph("Sector", s_label),
            Paragraph("Employees", s_label),
        ], [
            Paragraph(f'<font color="#{"%02x%02x%02x" % tuple(int(v*255) for v in _rgb(score_color.red, score_color.green, score_color.blue) if False) or "34C759"}">{health_score}/100 {health_label}</font>', s_h2)
            if False else Paragraph(f"{health_score}/100 — {health_label}", s_h2),
            Paragraph(_fmt_large(info.get("marketCap")), s_body),
            Paragraph(info.get("sector", "N/A"), s_body),
            Paragraph(f'{info.get("fullTimeEmployees",0):,}' if info.get("fullTimeEmployees") else "N/A", s_body),
        ]]
        score_tbl = Table(score_data, colWidths=[usable_w * 0.28] * 4)
        score_tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, -1), colors.Color(*_rgb(C_CARD))),
            ("GRID",       (0, 0), (-1, -1), 0.5, colors.Color(*_rgb(C_BORDER))),
            ("ROWBACKGROUNDS", (0, 0), (-1, -1), [colors.Color(*_rgb(C_CARD))] * 2),
            ("TOPPADDING",  (0, 0), (-1, -1), 8),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
            ("LEFTPADDING",  (0, 0), (-1, -1), 10),
        ]))
        story.append(score_tbl)
        story.append(Spacer(1, 10))

        # ── KPI table ─────────────────────────────────────────────────────────
        story.append(Paragraph("Key Performance Indicators", s_h1))
        kpi_display_map = [
            ("Revenue Growth %", "Revenue Growth"),
            ("Gross Margin %",   "Gross Margin"),
            ("Net Margin %",     "Net Margin"),
            ("Operating Margin %","Op Margin"),
            ("ROE %",            "ROE"),
            ("ROA %",            "ROA"),
            ("Current Ratio",    "Current Ratio"),
            ("Debt-to-Equity",   "Debt / Equity"),
            ("FCF Margin %",     "FCF Margin"),
        ]
        kpi_rows = [["Metric", "Value", "Metric", "Value"]]
        pairs = [(kpi_display_map[i], kpi_display_map[i+1])
                 for i in range(0, len(kpi_display_map)-1, 2)]
        for (k1, l1), (k2, l2) in pairs:
            v1 = kpis.get(k1, (None, "N/A"))[1]
            v2 = kpis.get(k2, (None, "N/A"))[1]
            kpi_rows.append([
                Paragraph(l1, s_label), Paragraph(v1, s_body),
                Paragraph(l2, s_label), Paragraph(v2, s_body),
            ])
        kpi_tbl = Table(kpi_rows, colWidths=[usable_w*0.25, usable_w*0.25]*2)
        kpi_tbl.setStyle(TableStyle([
            ("BACKGROUND",  (0, 0), (-1, 0),  colors.Color(*_rgb(C_BLUE))),
            ("TEXTCOLOR",   (0, 0), (-1, 0),  colors.white),
            ("FONTNAME",    (0, 0), (-1, 0),  "Helvetica-Bold"),
            ("FONTSIZE",    (0, 0), (-1, -1), 9),
            ("BACKGROUND",  (0, 1), (-1, -1), colors.Color(*_rgb(C_CARD))),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1),
             [colors.Color(*_rgb(C_CARD)), colors.Color(0.12, 0.12, 0.13)]),
            ("GRID",        (0, 0), (-1, -1), 0.5, colors.Color(*_rgb(C_BORDER))),
            ("TOPPADDING",  (0, 0), (-1, -1), 6),
            ("BOTTOMPADDING",(0, 0), (-1, -1), 6),
            ("LEFTPADDING", (0, 0), (-1, -1), 8),
        ]))
        story.append(kpi_tbl)
        story.append(Spacer(1, 10))

        # ── Executive Insights ────────────────────────────────────────────────
        if insights:
            story.append(Paragraph("Executive Insights", s_h1))
            icon_map = {
                "Revenue Trend": "Revenue",
                "Profitability Trend": "Profitability",
                "Balance Sheet Strength": "Balance Sheet",
                "Cash Flow Analysis": "Cash Flow",
            }
            for title, text in insights.items():
                story.append(Paragraph(icon_map.get(title, title), s_h2))
                story.append(Paragraph(text, s_body))
                story.append(Spacer(1, 4))

        # ── CFO Brief ─────────────────────────────────────────────────────────
        if cfo_brief:
            story.append(PageBreak())
            story.append(Paragraph("CFO Brief", s_h1))
            # Parse markdown headers into styled paragraphs
            for line in cfo_brief.split("\n"):
                line = line.strip()
                if not line:
                    story.append(Spacer(1, 4))
                elif line.startswith("## "):
                    story.append(Paragraph(line[3:], s_h2))
                elif line.startswith("**") and line.endswith("**"):
                    story.append(Paragraph(line[2:-2], s_h2))
                elif line.startswith("- "):
                    story.append(Paragraph(f"• {line[2:]}", s_body))
                elif line.startswith(str(tuple("123456789"))):
                    story.append(Paragraph(line, s_body))
                else:
                    story.append(Paragraph(line, s_body))

        # ── Recent News ────────────────────────────────────────────────────────
        if news_items:
            story.append(Spacer(1, 8))
            story.append(Paragraph("Recent News", s_h1))
            for item in news_items[:6]:
                story.append(Paragraph(f"• {item.get('title','')}", s_body))
                story.append(Paragraph(
                    f"  {item.get('publisher','')} · {item.get('published_at','')}",
                    s_muted,
                ))

        # ── Footer disclaimer ─────────────────────────────────────────────────
        story.append(Spacer(1, 16))
        story.append(HRFlowable(width="100%", thickness=0.5,
                                color=colors.Color(*_rgb(C_BORDER))))
        story.append(Spacer(1, 4))
        story.append(Paragraph(
            "Generated by FinIntel AI · Built by Hetal Shah · github.com/Hshah168 · "
            "Data via Yahoo Finance · Not financial advice.",
            s_muted,
        ))

        doc = SimpleDocTemplate(
            buf, pagesize=A4,
            leftMargin=margin, rightMargin=margin,
            topMargin=margin, bottomMargin=margin,
            title=f"{company_name} — FinIntel AI Report",
        )
        # Dark background on every page
        def dark_bg(canvas, doc):
            canvas.saveState()
            canvas.setFillColorRGB(*_rgb(C_BG))
            canvas.rect(0, 0, W, H, fill=True, stroke=False)
            canvas.restoreState()

        doc.build(story, onFirstPage=dark_bg, onLaterPages=dark_bg)
        return buf.getvalue()

    except Exception as e:
        # Return a simple fallback PDF with error message
        buf = io.BytesIO()
        try:
            from reportlab.pdfgen import canvas as rl_canvas
            c = rl_canvas.Canvas(buf)
            c.setFillColorRGB(1, 1, 1)
            c.drawString(50, 750, f"FinIntel AI Report — {company_name}")
            c.drawString(50, 730, f"PDF generation error: {str(e)}")
            c.drawString(50, 710, "Please download the CFO Brief as Markdown instead.")
            c.save()
        except Exception:
            pass
        return buf.getvalue()


def _fmt_large(n):
    if not n:
        return "N/A"
    try:
        n = float(n)
        if abs(n) >= 1e12: return f"${n/1e12:.2f}T"
        if abs(n) >= 1e9:  return f"${n/1e9:.2f}B"
        if abs(n) >= 1e6:  return f"${n/1e6:.2f}M"
        return f"${n:,.0f}"
    except Exception:
        return "N/A"


# ══════════════════════════════════════════════════════════════════════════════
# POWERPOINT EXPORT
# ══════════════════════════════════════════════════════════════════════════════

def generate_pptx(
    company_name: str,
    ticker: str,
    info: dict,
    kpis: dict,
    health_score: int,
    health_label: str,
    health_breakdown: dict,
    insights: dict,
    cfo_brief: str,
    price_data: dict,
) -> bytes:
    """Generate a 7-slide dark PowerPoint deck. Returns bytes."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        def rgb(c):
            return RGBColor(*c)

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]  # completely blank

        def add_slide():
            s = prs.slides.add_slide(blank_layout)
            bg = s.background.fill
            bg.solid()
            bg.fore_color.rgb = rgb(C_BG)
            return s

        def txb(slide, text, l, t, w, h,
                size=18, bold=False, color=C_WHITE,
                align=PP_ALIGN.LEFT, italic=False):
            tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = str(text)
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = rgb(color)
            return tb

        def rect(slide, l, t, w, h, fill=C_CARD, line=None):
            from pptx.util import Inches
            shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                Inches(l), Inches(t), Inches(w), Inches(h),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(fill)
            if line:
                shape.line.color.rgb = rgb(line)
                shape.line.width = Emu(9525)
            else:
                shape.line.fill.background()
            return shape

        def accent_bar(slide):
            rect(slide, 0, 0, 0.08, 7.5, fill=C_BLUE)

        today = datetime.date.today().strftime("%B %d, %Y")
        sector = info.get("sector", "N/A")
        mktcap = _fmt_large(info.get("marketCap"))
        price  = price_data.get("price", 0) or 0
        chg    = price_data.get("change_pct", 0) or 0

        score_color = (C_GREEN if health_score >= 75 else
                       C_BLUE  if health_score >= 55 else
                       C_ORANGE if health_score >= 35 else C_RED)

        # ── Slide 1: Cover ────────────────────────────────────────────────────
        s1 = add_slide()
        accent_bar(s1)
        rect(s1, 0.08, 0, 13.25, 7.5, fill=C_BG)
        # Blue accent top strip
        rect(s1, 0.08, 0, 13.25, 0.08, fill=C_BLUE)
        txb(s1, "FINANCIAL INTELLIGENCE REPORT", 0.5, 1.2, 12, 0.5,
            size=11, color=C_BLUE, bold=True)
        txb(s1, company_name, 0.5, 1.8, 12, 1.2, size=40, bold=True)
        txb(s1, f"{ticker}  ·  {sector}  ·  {today}",
            0.5, 3.1, 12, 0.5, size=14, color=C_MUTED)
        # Health score badge
        rect(s1, 0.5, 3.8, 3.5, 1.2, fill=C_CARD, line=C_BORDER)
        txb(s1, "FINANCIAL HEALTH", 0.6, 3.85, 3.3, 0.4,
            size=9, color=C_MUTED, bold=True)
        txb(s1, f"{health_score}/100  {health_label}",
            0.6, 4.2, 3.3, 0.6, size=20, bold=True, color=score_color)
        # Market data
        rect(s1, 4.2, 3.8, 3.0, 1.2, fill=C_CARD, line=C_BORDER)
        txb(s1, "MARKET CAP", 4.3, 3.85, 2.8, 0.4, size=9, color=C_MUTED, bold=True)
        txb(s1, mktcap, 4.3, 4.2, 2.8, 0.6, size=20, bold=True)

        rect(s1, 7.4, 3.8, 3.0, 1.2, fill=C_CARD, line=C_BORDER)
        txb(s1, "CURRENT PRICE", 7.5, 3.85, 2.8, 0.4, size=9, color=C_MUTED, bold=True)
        chg_c = C_GREEN if chg >= 0 else C_RED
        txb(s1, f"${price:.2f}", 7.5, 4.1, 2.8, 0.4, size=18, bold=True)
        txb(s1, f"{chg:+.2f}% today", 7.5, 4.55, 2.8, 0.35, size=11, color=chg_c)

        txb(s1, "Built by Hetal Shah · FinIntel AI · github.com/Hshah168",
            0.5, 6.9, 12, 0.4, size=8, color=C_MUTED,
            align=PP_ALIGN.CENTER)

        # ── Slide 2: KPI Dashboard ────────────────────────────────────────────
        s2 = add_slide()
        accent_bar(s2)
        txb(s2, "KEY PERFORMANCE INDICATORS", 0.3, 0.2, 12, 0.5,
            size=11, bold=True, color=C_BLUE)
        txb(s2, company_name, 0.3, 0.55, 12, 0.5, size=18, bold=True)

        kpi_list = [
            ("Revenue Growth", "Revenue Growth %"),
            ("Gross Margin",   "Gross Margin %"),
            ("Net Margin",     "Net Margin %"),
            ("Op Margin",      "Operating Margin %"),
            ("ROE",            "ROE %"),
            ("ROA",            "ROA %"),
            ("Current Ratio",  "Current Ratio"),
            ("Debt / Equity",  "Debt-to-Equity"),
            ("FCF Margin",     "FCF Margin %"),
        ]
        cols, rows_per_col = 3, 3
        col_w, row_h = 4.2, 1.65
        start_x, start_y = 0.2, 1.1

        for i, (label, key) in enumerate(kpi_list):
            col_i = i % cols
            row_i = i // cols
            x = start_x + col_i * col_w
            y = start_y + row_i * row_h
            val = kpis.get(key, (None, "N/A"))[1]
            raw = kpis.get(key, (None,))[0]
            vcolor = C_WHITE
            if raw is not None:
                if key in ["Revenue Growth %", "Gross Margin %", "Net Margin %",
                           "ROE %", "ROA %", "FCF Margin %", "Operating Margin %"]:
                    vcolor = C_GREEN if raw > 0 else C_RED
                elif key == "Current Ratio":
                    vcolor = C_GREEN if raw >= 1.5 else C_ORANGE if raw >= 1.0 else C_RED
                elif key == "Debt-to-Equity":
                    vcolor = C_GREEN if raw < 0.5 else C_ORANGE if raw < 1.5 else C_RED
            rect(s2, x, y, col_w - 0.15, row_h - 0.15, fill=C_CARD, line=C_BORDER)
            txb(s2, label.upper(), x + 0.15, y + 0.1, col_w - 0.3, 0.35,
                size=8, color=C_MUTED, bold=True)
            txb(s2, val, x + 0.15, y + 0.45, col_w - 0.3, 0.7,
                size=24, bold=True, color=vcolor)

        # ── Slide 3: Financial Health Score ───────────────────────────────────
        s3 = add_slide()
        accent_bar(s3)
        txb(s3, "FINANCIAL HEALTH SCORE", 0.3, 0.2, 12, 0.5,
            size=11, bold=True, color=C_BLUE)
        txb(s3, company_name, 0.3, 0.55, 12, 0.5, size=18, bold=True)

        # Big score
        rect(s3, 0.3, 1.1, 4.0, 3.5, fill=C_CARD, line=C_BORDER)
        txb(s3, str(health_score), 0.5, 1.5, 3.6, 1.8,
            size=72, bold=True, color=score_color, align=PP_ALIGN.CENTER)
        txb(s3, "out of 100", 0.5, 3.1, 3.6, 0.4,
            size=12, color=C_MUTED, align=PP_ALIGN.CENTER)
        txb(s3, health_label.upper(), 0.5, 3.5, 3.6, 0.6,
            size=16, bold=True, color=score_color, align=PP_ALIGN.CENTER)

        # Dimension bars
        dims = list(health_breakdown.items())
        for i, (dim, score) in enumerate(dims):
            y = 1.1 + i * 0.62
            pct = score / 20
            dc = (C_GREEN if pct >= 0.75 else C_BLUE if pct >= 0.5
                  else C_ORANGE if pct >= 0.25 else C_RED)
            rect(s3, 4.6, y, 8.0, 0.52, fill=C_CARD, line=C_BORDER)
            txb(s3, dim, 4.75, y + 0.05, 4.0, 0.3, size=10, color=C_WHITE)
            txb(s3, f"{score}/20", 11.8, y + 0.05, 0.7, 0.3,
                size=10, bold=True, color=dc, align=PP_ALIGN.RIGHT)
            # Progress bar background
            rect(s3, 4.75, y + 0.32, 7.5, 0.14, fill=C_BORDER)
            # Progress bar fill
            if pct > 0:
                rect(s3, 4.75, y + 0.32, max(0.1, 7.5 * pct), 0.14, fill=dc)

        # ── Slide 4: Executive Insights ───────────────────────────────────────
        s4 = add_slide()
        accent_bar(s4)
        txb(s4, "EXECUTIVE INSIGHTS", 0.3, 0.2, 12, 0.5,
            size=11, bold=True, color=C_BLUE)
        txb(s4, company_name, 0.3, 0.55, 12, 0.5, size=18, bold=True)

        insight_items = list(insights.items())[:4]
        positions = [(0.3, 1.15), (6.8, 1.15), (0.3, 4.05), (6.8, 4.05)]
        for (title, text), (x, y) in zip(insight_items, positions):
            rect(s4, x, y, 6.3, 2.75, fill=C_CARD, line=C_BORDER)
            rect(s4, x, y, 6.3, 0.08, fill=C_BLUE)
            txb(s4, title.upper(), x + 0.15, y + 0.12, 6.0, 0.35,
                size=9, bold=True, color=C_BLUE)
            # Truncate text to fit
            short = text[:280] + "..." if len(text) > 280 else text
            txb(s4, short, x + 0.15, y + 0.5, 6.0, 2.1, size=9, color=C_WHITE)

        # ── Slide 5: CFO Brief Summary ────────────────────────────────────────
        s5 = add_slide()
        accent_bar(s5)
        txb(s5, "CFO BRIEF", 0.3, 0.2, 12, 0.5, size=11, bold=True, color=C_BLUE)
        txb(s5, company_name, 0.3, 0.55, 12, 0.5, size=18, bold=True)

        # Parse CFO brief into sections
        sections = []
        curr_section = None
        curr_lines   = []
        for line in (cfo_brief or "").split("\n"):
            line = line.strip()
            if line.startswith("## "):
                if curr_section:
                    sections.append((curr_section, " ".join(curr_lines)))
                curr_section = line[3:]
                curr_lines = []
            elif line and not line.startswith("#"):
                curr_lines.append(line)
        if curr_section:
            sections.append((curr_section, " ".join(curr_lines)))

        col_positions = [(0.3, 1.1), (6.8, 1.1), (0.3, 3.9), (6.8, 3.9),
                         (0.3, 5.8), (6.8, 5.8)]
        for (title, body), (x, y) in zip(sections[:6], col_positions):
            h = 2.6 if y < 3 else (1.5 if y > 5 else 2.6)
            rect(s5, x, y, 6.3, h, fill=C_CARD, line=C_BORDER)
            txb(s5, title.upper(), x + 0.15, y + 0.1, 6.0, 0.35,
                size=9, bold=True, color=C_BLUE)
            short = body[:260] + "..." if len(body) > 260 else body
            txb(s5, short, x + 0.15, y + 0.45, 6.0, h - 0.55, size=9, color=C_WHITE)

        # ── Slide 6: Balance Sheet Snapshot ───────────────────────────────────
        s6 = add_slide()
        accent_bar(s6)
        txb(s6, "BALANCE SHEET SNAPSHOT", 0.3, 0.2, 12, 0.5,
            size=11, bold=True, color=C_BLUE)
        txb(s6, company_name, 0.3, 0.55, 12, 0.5, size=18, bold=True)

        bs_metrics = [
            ("52W High",      f"${price_data.get('high_52w') or 0:.2f}"),
            ("52W Low",       f"${price_data.get('low_52w') or 0:.2f}"),
            ("Beta",          f"{price_data.get('beta') or 'N/A'}"),
            ("P/E Ratio",     f"{price_data.get('pe_ratio') or 'N/A'}"),
            ("Div Yield",     f"{(price_data.get('dividend_yield') or 0)*100:.2f}%"),
            ("EPS",           f"${price_data.get('eps') or 'N/A'}"),
        ]
        for i, (label, val) in enumerate(bs_metrics):
            x = 0.3 + (i % 3) * 4.3
            y = 1.4 + (i // 3) * 1.6
            rect(s6, x, y, 4.1, 1.4, fill=C_CARD, line=C_BORDER)
            txb(s6, label.upper(), x + 0.15, y + 0.1, 3.8, 0.35,
                size=9, color=C_MUTED, bold=True)
            txb(s6, val, x + 0.15, y + 0.5, 3.8, 0.7, size=22, bold=True)

        txb(s6, info.get("longBusinessSummary", "")[:400],
            0.3, 4.8, 12.7, 2.3, size=9, color=C_MUTED)

        # ── Slide 7: Closing ──────────────────────────────────────────────────
        s7 = add_slide()
        rect(s7, 0, 0, 13.33, 7.5, fill=C_BG)
        rect(s7, 0, 0, 13.33, 0.08, fill=C_BLUE)
        rect(s7, 0, 7.42, 13.33, 0.08, fill=C_BLUE)
        txb(s7, "FinIntel AI", 0.5, 2.5, 12.3, 1.0,
            size=36, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)
        txb(s7, "Financial Intelligence Platform", 0.5, 3.4, 12.3, 0.6,
            size=16, color=C_MUTED, align=PP_ALIGN.CENTER)
        txb(s7, "Built by Hetal Shah · github.com/Hshah168", 0.5, 4.2, 12.3, 0.5,
            size=12, color=C_MUTED, align=PP_ALIGN.CENTER)
        txb(s7, "Data via Yahoo Finance · Not financial advice · Generated by FinIntel AI",
            0.5, 6.8, 12.3, 0.4, size=8, color=C_MUTED, align=PP_ALIGN.CENTER)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    except Exception as e:
        # Return minimal PPTX on error
        try:
            from pptx import Presentation
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = f"FinIntel AI — {company_name}"
            slide.placeholders[1].text = f"Export error: {str(e)}"
            buf = io.BytesIO()
            prs.save(buf)
            return buf.getvalue()
        except Exception:
            return b""
